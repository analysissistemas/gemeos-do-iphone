"""
Gera a apresentação em PowerPoint (.pptx).

Rodar:  python gerar_ppt.py
Saída:  Gemeos-do-iPhone-Sistema.pptx

Paleta: preto e branco da marca, com o verde do WhatsApp como único acento.
O verde não é enfeite — o WhatsApp é o canal de venda da loja, então a cor
marca exatamente as ações que geram dinheiro.
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

AQUI = os.path.dirname(os.path.abspath(__file__))
img = lambda n: os.path.join(AQUI, n)

PRETO = RGBColor(0x0B, 0x0B, 0x0C)
CARVAO = RGBColor(0x16, 0x16, 0x1A)
CLARO = RGBColor(0xF5, 0xF5, 0xF7)
PAPEL = RGBColor(0xFF, 0xFF, 0xFF)
NEVE = RGBColor(0xF4, 0xF6, 0xFA)
TINTA = RGBColor(0x10, 0x18, 0x28)
CINZA = RGBColor(0xA1, 0xA1, 0xA6)
CINZA_C = RGBColor(0x5B, 0x64, 0x72)
FRACO = RGBColor(0x4A, 0x4A, 0x52)
BORDA_E = RGBColor(0x2A, 0x2A, 0x30)
BORDA_C = RGBColor(0xE2, 0xE6, 0xEE)
VERDE = RGBColor(0x25, 0xD3, 0x66)
VERDE_ESC = RGBColor(0x05, 0x38, 0x1A)

SERIF = "Cambria"      # títulos — da lista segura, renderiza igual em todo Office
SANS = "Calibri"       # corpo

L = Inches(0.7)        # margem esquerda
LARGA = Inches(13.333)
ALTA = Inches(7.5)

pres = Presentation()
pres.slide_width = LARGA
pres.slide_height = ALTA
VAZIO = pres.slide_layouts[6]      # layout em branco


def slide(cor_fundo):
    s = pres.slides.add_slide(VAZIO)
    fundo = s.background.fill
    fundo.solid()
    fundo.fore_color.rgb = cor_fundo
    return s


def texto(s, txt, x, y, w, h, *, tam=12, cor=CINZA, negrito=False, fonte=SANS,
          alinha=PP_ALIGN.LEFT, espaco=None, italico=False, espacamento=None):
    cx = s.shapes.add_textbox(x, y, w, h)
    tf = cx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, linha in enumerate(str(txt).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alinha
        if espacamento:
            p.line_spacing = espacamento
        r = p.add_run()
        r.text = linha
        f = r.font
        f.size = Pt(tam)
        f.bold = negrito
        f.italic = italico
        f.color.rgb = cor
        f.name = fonte
        if espaco:                      # espaçamento entre letras, em pontos
            r.font._rPr.set("spc", str(int(espaco * 100)))
    return cx


def caixa(s, x, y, w, h, *, preenche, borda=None, raio=True):
    forma = MSO_SHAPE.ROUNDED_RECTANGLE if raio else MSO_SHAPE.RECTANGLE
    sh = s.shapes.add_shape(forma, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = preenche
    if borda:
        sh.line.color.rgb = borda
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    if raio:                            # canto discreto, não bolha
        sh.adjustments[0] = 0.04
    return sh


def bolinha(s, n, x, y, *, fundo, tinta_num, d=Inches(0.34)):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fundo
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(n)
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = tinta_num
    r.font.name = SANS


def etiqueta(s, txt, escuro=True):
    texto(s, txt.upper(), L, Inches(0.55), Inches(6), Inches(0.3),
          tam=10, negrito=True, espaco=2.2, cor=CINZA if escuro else CINZA_C)


def titulo(s, txt, escuro=True, y=0.92):
    texto(s, txt, L, Inches(y), Inches(11.5), Inches(0.85),
          tam=34, negrito=True, fonte=SERIF, cor=CLARO if escuro else TINTA)


def linha_fina(s, txt, escuro=True, y=1.80, w=11.0):
    texto(s, txt, L, Inches(y), Inches(w), Inches(0.62),
          tam=13, cor=CINZA if escuro else CINZA_C, espacamento=1.35)


def rodape(s, n, escuro=True):
    cor = FRACO if escuro else RGBColor(0x98, 0xA2, 0xB3)
    texto(s, "Gêmeos do iPhone", L, Inches(6.95), Inches(5), Inches(0.3), tam=9, cor=cor)
    texto(s, str(n), Inches(12.0), Inches(6.95), Inches(0.63), Inches(0.3),
          tam=9, cor=cor, alinha=PP_ALIGN.RIGHT)


def bloco(s, n, tit, txt, x, y, w, escuro=True):
    """O motivo visual que se repete: número em círculo + título + explicação."""
    bolinha(s, n, x, y, fundo=CLARO if escuro else TINTA,
            tinta_num=PRETO if escuro else PAPEL)
    texto(s, tit, x + Inches(0.48), y - Inches(0.02), w - Inches(0.48), Inches(0.28),
          tam=13, negrito=True, cor=CLARO if escuro else TINTA)
    texto(s, txt, x + Inches(0.48), y + Inches(0.27), w - Inches(0.48), Inches(0.62),
          tam=11, cor=CINZA if escuro else CINZA_C, espacamento=1.2)


def tela(s, arquivo, x, y, w, h, escuro=True):
    caixa(s, x - Inches(0.05), y - Inches(0.05), w + Inches(0.1), h + Inches(0.1),
          preenche=CARVAO if escuro else NEVE, borda=BORDA_E if escuro else BORDA_C)
    s.shapes.add_picture(img(arquivo), x, y, w, h)


def blocos(s, lista, escuro=True, x=8.35, y0=2.60, dy=1.03, w=4.3):
    for i, (n, t, d) in enumerate(lista):
        bloco(s, n, t, d, Inches(x), Inches(y0 + i * dy), Inches(w), escuro)


# ============================================================ 1. CAPA
s = slide(PRETO)
halo = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.6), Inches(-2.4), Inches(8.1), Inches(8.1))
halo.fill.solid()
halo.fill.fore_color.rgb = RGBColor(0x18, 0x18, 0x1D)
halo.line.fill.background()
halo.shadow.inherit = False
s.shapes.add_picture(img("logo.png"), Inches(5.15), Inches(1.35), Inches(3.0), Inches(0.98))
texto(s, "O sistema da loja,", Inches(0.9), Inches(2.70), Inches(11.5), Inches(0.78),
      tam=44, negrito=True, fonte=SERIF, cor=CLARO, alinha=PP_ALIGN.CENTER)
texto(s, "de ponta a ponta.", Inches(0.9), Inches(3.48), Inches(11.5), Inches(0.78),
      tam=44, negrito=True, fonte=SERIF, cor=RGBColor(0x8E, 0x8E, 0x93), alinha=PP_ALIGN.CENTER)
texto(s, "Da vitrine que o cliente vê ao controle que só a equipe enxerga.\n"
         "Estoque, vendas, assistência técnica e caixa em um lugar só.",
      Inches(2.4), Inches(4.45), Inches(8.5), Inches(0.8),
      tam=14, cor=CINZA, alinha=PP_ALIGN.CENTER, espacamento=1.4)
sel = caixa(s, Inches(4.75), Inches(5.50), Inches(3.8), Inches(0.42),
            preenche=PRETO, borda=RGBColor(0x3A, 0x3A, 0x42))
texto(s, "CARPINA · GOIANA — PERNAMBUCO", Inches(4.75), Inches(5.60), Inches(3.8), Inches(0.3),
      tam=9.5, cor=CINZA, alinha=PP_ALIGN.CENTER, espaco=1.4)
texto(s, "Apresentação · 15 de agosto de 2026", Inches(0.9), Inches(6.60), Inches(11.5),
      Inches(0.3), tam=10, cor=FRACO, alinha=PP_ALIGN.CENTER)
s.notes_slide.notes_text_frame.text = (
    "Abrir dizendo: o sistema tem dois lados — o que o cliente vê e o que só a "
    "equipe vê. Os dois se alimentam do mesmo cadastro.")

# ============================================================ 2. DOIS LADOS
s = slide(PRETO)
etiqueta(s, "Como funciona")
titulo(s, "Dois lados, um cadastro só.")
linha_fina(s, "Você dá entrada no aparelho uma vez. Ele aparece na loja para o cliente "
              "e no controle para a equipe — sem digitar nada duas vezes.")
LADOS = [
    (0.70, "O CLIENTE VÊ", "A loja", VERDE,
     ["Foto do aparelho, que troca ao escolher a cor",
      "Memória, condição e saúde da bateria",
      "Avarias declaradas abertamente",
      "Preço e parcela em até 18x",
      "Botão que abre o WhatsApp com a mensagem pronta"],
     "Não vê: quanto você pagou, sua margem, nem o lucro."),
    (6.85, "A EQUIPE VÊ", "O sistema", CLARO,
     ["Custo, preço de venda, margem e lucro por peça",
      "IMEI, bateria e avarias de cada aparelho",
      "Funil de vendas com cada negócio em aberto",
      "Ordens de serviço da assistência técnica",
      "Caixa: recebido, a receber e parcelas atrasadas"],
     "Protegido por login, com acesso separado por pessoa."),
]
for x, quem, nome, cor, itens, nao in LADOS:
    caixa(s, Inches(x), Inches(2.62), Inches(5.78), Inches(3.95),
          preenche=CARVAO, borda=BORDA_E)
    texto(s, quem, Inches(x + 0.4), Inches(2.92), Inches(4.9), Inches(0.25),
          tam=9.5, negrito=True, cor=cor, espaco=1.8)
    texto(s, nome, Inches(x + 0.4), Inches(3.20), Inches(4.9), Inches(0.45),
          tam=22, negrito=True, fonte=SERIF, cor=CLARO)
    for i, it in enumerate(itens):
        y = 3.82 + i * 0.42
        pt = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.42), Inches(y + 0.07),
                                Inches(0.07), Inches(0.07))
        pt.fill.solid()
        pt.fill.fore_color.rgb = RGBColor(0x6E, 0x6E, 0x73)
        pt.line.fill.background()
        pt.shadow.inherit = False
        texto(s, it, Inches(x + 0.62), Inches(y), Inches(4.85), Inches(0.36),
              tam=11.5, cor=CINZA)
    texto(s, nao, Inches(x + 0.4), Inches(5.95), Inches(4.95), Inches(0.45),
          tam=10.5, cor=RGBColor(0x76, 0x76, 0x7E), italico=True)
rodape(s, 2)
s.notes_slide.notes_text_frame.text = (
    "Ponto forte de venda: o custo e o lucro não existem dentro do arquivo do "
    "cliente. Não é só esconder na tela — quem abrir o código-fonte não acha.")

# ============================================================ 3. A LOJA
s = slide(PRETO)
etiqueta(s, "Parte do cliente")
titulo(s, "A loja aberta 24 horas.")
linha_fina(s, "O cliente entra pelo link, vê o estoque real e chama no WhatsApp já "
              "sabendo o que quer.")
tela(s, "01-loja-topo.jpg", L, Inches(2.55), Inches(7.3), Inches(3.64))
blocos(s, [
    (1, "Identidade da loja", "Logo, cores e a frase da casa: \"Pra tudo tem solução\"."),
    (2, "Vídeo de fundo", "Abre com movimento, como as grandes lojas de tecnologia."),
    (3, "Funciona no celular", "A maioria dos clientes vai abrir pelo telefone."),
    (4, "WhatsApp sempre à mão", "Botão fixo na tela, em qualquer página."),
])
rodape(s, 3)

# ============================================================ 4. CATÁLOGO
s = slide(PAPEL)
etiqueta(s, "Parte do cliente", False)
titulo(s, "Catálogo que responde sozinho.", False)
linha_fina(s, "Seis categorias — iPhone, iPad, Apple Watch, MacBook, acessórios e motos "
              "elétricas — com busca por modelo, cor ou memória.", False)
tela(s, "02-catalogo.jpg", L, Inches(2.55), Inches(7.3), Inches(3.64), escuro=False)
blocos(s, [
    (1, "Escolhe a cor, a foto muda", "Cores reais de cada geração, com a foto oficial do aparelho."),
    (2, "Bateria e avarias à mostra", "O que os outros escondem, aqui está escrito no anúncio."),
    (3, "Preço e parcela na hora", "Valor à vista e em 18x, sem o cliente precisar perguntar."),
    (4, "Mensagem pronta", "Chega com modelo, cor, memória, bateria e preço."),
], escuro=False)
rodape(s, 4, False)

# ============================================================ 5. ESTOQUE
s = slide(PRETO)
etiqueta(s, "Parte da equipe")
titulo(s, "Cada aparelho é uma peça única.")
linha_fina(s, "Bateria e avaria são daquele aparelho, não do modelo. Por isso o controle "
              "é peça por peça, com IMEI — nunca por quantidade.")
tela(s, "03-estoque.jpg", L, Inches(2.55), Inches(7.3), Inches(3.64))
blocos(s, [
    (1, "Lucro por aparelho", "Custo, venda e lucro lado a lado, na mesma linha."),
    (2, "Preço sugerido", "Calculado pela condição, bateria e avarias — mostrando a conta."),
    (3, "Filtros que importam", "Com avaria, bateria baixa, lacrado, seminovo, vitrine."),
    (4, "Cadastro em segundos", "Escolhido o modelo, só oferece memórias e cores que existem."),
])
rodape(s, 5)
s.notes_slide.notes_text_frame.text = (
    "Explicar por que peça única: um iPhone 13 com 87% de bateria e outro com 100% "
    "não valem o mesmo. Contar por quantidade esconde isso.")

# ============================================================ 6. FUNIL
s = slide(PAPEL)
etiqueta(s, "Parte da equipe", False)
titulo(s, "Nenhuma venda esquecida.", False)
linha_fina(s, "Cada conversa vira um cartão, da primeira mensagem até o fechamento — e dá "
              "para ver quanto dinheiro está em jogo agora.", False)
tela(s, "04-funil.jpg", L, Inches(2.55), Inches(7.3), Inches(3.64), escuro=False)
blocos(s, [
    (1, "Cinco etapas", "Chegou no WhatsApp, qualificado, proposta, troca, fechada."),
    (2, "Conversão por etapa", "Mostra onde o cliente desiste, para atacar o ponto certo."),
    (3, "Valor em jogo", "Quanto está parado no funil, em reais, agora."),
    (4, "Alerta de parado", "Negócio sem resposta há mais de 7 dias aparece na Visão Geral."),
], escuro=False)
rodape(s, 6, False)

# ============================================================ 7. ASSISTÊNCIA
s = slide(PRETO)
etiqueta(s, "Parte da equipe")
titulo(s, "A bancada sob controle.")
linha_fina(s, "Toda ordem de serviço, da entrada até a entrega, com prazo e técnico "
              "responsável — e a nota do cliente sai pronta para imprimir.")
tela(s, "05-assistencia.jpg", L, Inches(2.55), Inches(7.3), Inches(3.64))
blocos(s, [
    (1, "Seis etapas", "Recebido, orçamento, aprovado, bancada, aguardando peça, pronto."),
    (2, "Aparelho parado aparece", "Quem está esperando peça, e há quantos dias."),
    (3, "Nota com garantia", "Defeito, valor, prazo de garantia e as duas assinaturas."),
    (4, "Quanto vai entrar", "Soma dos orçamentos aprovados que ainda serão recebidos."),
])
rodape(s, 7)

# ============================================================ 8. NÚMEROS
s = slide(PAPEL)
etiqueta(s, "Onde estamos", False)
titulo(s, "O que já está de pé.", False)
linha_fina(s, "Tudo abaixo funciona hoje e pode ser aberto agora, no celular de qualquer "
              "pessoa da equipe.", False)
for i, (n, t, d) in enumerate([
        ("8", "telas prontas", "entre a loja e o sistema"),
        ("6", "categorias", "no catálogo do cliente"),
        ("78", "fotos oficiais", "de produto, uma por cor")]):
    x = 0.70 + i * 3.98
    caixa(s, Inches(x), Inches(2.75), Inches(3.62), Inches(1.85),
          preenche=NEVE, borda=BORDA_C)
    texto(s, n, Inches(x + 0.35), Inches(2.92), Inches(2.9), Inches(0.9),
          tam=46, negrito=True, fonte=SERIF, cor=TINTA)
    texto(s, t, Inches(x + 0.35), Inches(3.84), Inches(2.9), Inches(0.28),
          tam=13, negrito=True, cor=TINTA)
    texto(s, d, Inches(x + 0.35), Inches(4.11), Inches(2.9), Inches(0.3),
          tam=11, cor=CINZA_C)
caixa(s, L, Inches(4.95), Inches(11.9), Inches(1.38), preenche=PRETO)
texto(s, "A vantagem que nenhum concorrente da região tem",
      Inches(1.15), Inches(5.18), Inches(11), Inches(0.32),
      tam=14, negrito=True, cor=VERDE)
texto(s, "Loja de seminovo vive de confiança. O sistema publica a saúde da bateria e as "
         "avarias de cada aparelho no próprio anúncio — o cliente vê o estado real antes "
         "de perguntar. Quem esconde não consegue competir com isso.",
      Inches(1.15), Inches(5.52), Inches(10.9), Inches(0.65),
      tam=11.5, cor=CINZA, espacamento=1.25)
rodape(s, 8, False)

# ============================================================ 9. PRÓXIMOS PASSOS
s = slide(PRETO)
etiqueta(s, "O caminho")
titulo(s, "Próximos passos.")
linha_fina(s, "Em ordem de impacto: cada um destrava o seguinte.")
for i, (t, d) in enumerate([
        ("Foto pelo celular, no cadastro",
         "Fotografar o aparelho na entrada resolve todo modelo sem foto oficial — e mostra "
         "ao cliente a peça real que ele vai receber."),
        ("Guardar os dados de verdade",
         "Hoje o cadastro vive no navegador de cada máquina. Com servidor, passa a valer "
         "para toda a equipe ao mesmo tempo."),
        ("Endereço próprio na internet",
         "Um gemeosdoiphone.com.br no lugar do link interno, para o cliente chegar de "
         "qualquer lugar."),
        ("Integração com o WhatsApp",
         "Conectar ao atendimento para o cliente sair da loja direto na conversa, já "
         "identificado.")]):
    y = 2.60 + i * 1.08
    caixa(s, L, Inches(y), Inches(11.9), Inches(0.94), preenche=CARVAO, borda=BORDA_E)
    bolinha(s, i + 1, Inches(1.05), Inches(y + 0.29), fundo=VERDE, tinta_num=VERDE_ESC,
            d=Inches(0.36))
    texto(s, t, Inches(1.60), Inches(y + 0.17), Inches(10.6), Inches(0.3),
          tam=14, negrito=True, cor=CLARO)
    texto(s, d, Inches(1.60), Inches(y + 0.48), Inches(10.5), Inches(0.36),
          tam=11, cor=CINZA)
rodape(s, 9)
s.notes_slide.notes_text_frame.text = (
    "Fechar no item 1: é o que resolve o problema das fotos de vez e não depende "
    "de ninguém de fora.")

saida = os.path.join(AQUI, "Gemeos-do-iPhone-Sistema.pptx")
pres.save(saida)
print("gerado:", saida)
print("slides:", len(pres.slides.__iter__.__self__._sldIdLst))
